"""
Document processing service for extracting text, images, and metadata
"""
import base64
import io
import os
from typing import List, Tuple, Dict, Any, Optional
import fitz  # PyMuPDF
from PIL import Image
import cv2
import numpy as np
from pathlib import Path

from app.models.schemas import DocumentMetadata


class DocumentProcessor:
    """Handles document parsing, image extraction, and pre-processing checks"""

    SUPPORTED_FORMATS = {
        '.pdf',  # PDF documents
        '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif', '.webp',  # Images
        '.docx', '.doc',  # Word documents
        '.pptx', '.ppt',  # PowerPoint
        '.xlsx', '.xls',  # Excel
        '.txt', '.md', '.csv',  # Text files
    }
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    MIN_LEGIBILITY_SCORE = 100.0  # Laplacian variance threshold

    def __init__(self):
        pass

    def _detect_image_format(self, image_bytes: bytes) -> str:
        """
        Detect image format from bytes (magic numbers)
        Returns MIME type for Claude API
        """
        # Check magic numbers (first few bytes)
        if len(image_bytes) < 12:
            return "image/png"  # Default fallback

        # PNG: 89 50 4E 47
        if image_bytes[:4] == b'\x89PNG':
            return "image/png"

        # JPEG: FF D8 FF
        if image_bytes[:3] == b'\xFF\xD8\xFF':
            return "image/jpeg"

        # GIF: 47 49 46
        if image_bytes[:3] == b'GIF':
            return "image/gif"

        # WebP: RIFF ... WEBP
        if image_bytes[:4] == b'RIFF' and image_bytes[8:12] == b'WEBP':
            return "image/webp"

        # BMP: 42 4D
        if image_bytes[:2] == b'BM':
            return "image/bmp"

        # TIFF: 49 49 or 4D 4D
        if image_bytes[:2] in (b'II', b'MM'):
            return "image/tiff"

        # Default to PNG if unknown
        return "image/png"

    def process_document(self, file_path: str) -> Tuple[DocumentMetadata, List[str], List[str]]:
        """
        Process a document and extract metadata, text, and images

        Returns:
            Tuple of (metadata, text_pages, image_base64_list)
        """
        file_ext = Path(file_path).suffix.lower()

        if file_ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported file format: {file_ext}")

        file_size = os.path.getsize(file_path)
        if file_size > self.MAX_FILE_SIZE:
            raise ValueError(f"File too large: {file_size} bytes (max {self.MAX_FILE_SIZE})")

        # Route to appropriate processor
        if file_ext == '.pdf':
            return self._process_pdf(file_path)
        elif file_ext in {'.docx', '.doc'}:
            return self._process_word(file_path)
        elif file_ext in {'.pptx', '.ppt'}:
            return self._process_powerpoint(file_path)
        elif file_ext in {'.xlsx', '.xls'}:
            return self._process_excel(file_path)
        elif file_ext in {'.txt', '.md', '.csv'}:
            return self._process_text(file_path)
        else:
            return self._process_image(file_path)

    def _process_pdf(self, file_path: str) -> Tuple[DocumentMetadata, List[str], List[str]]:
        """Process PDF document"""
        warnings = []
        text_pages = []
        images_base64 = []

        try:
            doc = fitz.open(file_path)
            page_count = len(doc)
            image_count = 0
            legibility_scores = []

            for page_num, page in enumerate(doc, 1):
                # Extract text
                text = page.get_text()
                text_pages.append(text)

                # Extract images from page
                image_list = page.get_images(full=True)
                image_count += len(image_list)

                # Extract each image
                for img_index, img_info in enumerate(image_list):
                    try:
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]

                        # Convert to base64
                        img_base64 = base64.b64encode(image_bytes).decode('utf-8')
                        images_base64.append(img_base64)

                        # Check image legibility
                        img = Image.open(io.BytesIO(image_bytes))
                        score = self._check_image_legibility(img)
                        legibility_scores.append(score)

                    except Exception as e:
                        warnings.append(f"Failed to extract image {img_index} from page {page_num}: {str(e)}")

                # Also render page as image for visual analysis
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x scale for better quality
                    img_bytes = pix.tobytes("png")
                    img_base64 = base64.b64encode(img_bytes).decode('utf-8')
                    images_base64.append(img_base64)

                    # Check page rendering legibility
                    img = Image.open(io.BytesIO(img_bytes))
                    score = self._check_image_legibility(img)
                    legibility_scores.append(score)

                except Exception as e:
                    warnings.append(f"Failed to render page {page_num} as image: {str(e)}")

            doc.close()

            # Determine overall legibility
            avg_legibility = sum(legibility_scores) / len(legibility_scores) if legibility_scores else 0
            is_legible = avg_legibility >= self.MIN_LEGIBILITY_SCORE

            if not is_legible:
                warnings.append(f"Low legibility detected (score: {avg_legibility:.2f})")

            if page_count == 0:
                warnings.append("Document has no pages")

            metadata = DocumentMetadata(
                filename=Path(file_path).name,
                file_size=os.path.getsize(file_path),
                page_count=page_count,
                image_count=len(images_base64),
                is_legible=is_legible,
                legibility_score=avg_legibility,
                format="PDF",
                warnings=warnings
            )

            return metadata, text_pages, images_base64

        except Exception as e:
            raise ValueError(f"Error processing PDF: {str(e)}")

    def _process_image(self, file_path: str) -> Tuple[DocumentMetadata, List[str], List[str]]:
        """Process standalone image file"""
        warnings = []

        try:
            with open(file_path, 'rb') as f:
                image_bytes = f.read()

            img = Image.open(file_path)

            # Check legibility
            legibility_score = self._check_image_legibility(img)
            is_legible = legibility_score >= self.MIN_LEGIBILITY_SCORE

            if not is_legible:
                warnings.append(f"Low image quality detected (score: {legibility_score:.2f})")

            # Convert to base64
            img_base64 = base64.b64encode(image_bytes).decode('utf-8')

            metadata = DocumentMetadata(
                filename=Path(file_path).name,
                file_size=os.path.getsize(file_path),
                page_count=1,
                image_count=1,
                is_legible=is_legible,
                legibility_score=legibility_score,
                format=Path(file_path).suffix.upper().replace('.', ''),
                warnings=warnings
            )

            # No text extraction for images (Claude will handle OCR)
            return metadata, [], [img_base64]

        except Exception as e:
            raise ValueError(f"Error processing image: {str(e)}")

    def _process_word(self, file_path: str) -> Tuple[DocumentMetadata, List[str], List[str]]:
        """Process Word documents (.docx, .doc)"""
        warnings = []
        text_pages = []
        images_base64 = []

        try:
            from docx import Document
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import _Cell, Table
            from docx.text.paragraph import Paragraph
            
            doc = Document(file_path)
            
            # Extract text by paragraphs (simulate pages by breaking every ~500 words)
            current_page = []
            word_count = 0
            
            for element in doc.element.body:
                if isinstance(element, CT_P):
                    para = Paragraph(element, doc)
                    text = para.text.strip()
                    if text:
                        current_page.append(text)
                        word_count += len(text.split())
                        
                        # Break into "pages" every ~500 words
                        if word_count >= 500:
                            text_pages.append('\n'.join(current_page))
                            current_page = []
                            word_count = 0
                            
                elif isinstance(element, CT_Tbl):
                    table = Table(element, doc)
                    table_text = []
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        table_text.append(' | '.join(row_text))
                    current_page.append('\n'.join(table_text))
            
            # Add remaining text
            if current_page:
                text_pages.append('\n'.join(current_page))
            
            # Ensure at least one page
            if not text_pages:
                text_pages = ["[No text content found]"]
            
            # Extract images
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        img_base64 = base64.b64encode(image_data).decode('utf-8')
                        images_base64.append(img_base64)
                    except Exception:
                        pass
            
            metadata = DocumentMetadata(
                filename=Path(file_path).name,
                file_size=os.path.getsize(file_path),
                page_count=len(text_pages),
                image_count=len(images_base64),
                is_legible=True,
                legibility_score=999.0,  # Text documents are always legible
                format="DOCX" if file_path.endswith('.docx') else "DOC",
                warnings=warnings
            )
            
            return metadata, text_pages, images_base64
            
        except ImportError:
            raise ValueError("python-docx not installed. Run: pip install python-docx")
        except Exception as e:
            raise ValueError(f"Error processing Word document: {str(e)}")

    def _process_powerpoint(self, file_path: str) -> Tuple[DocumentMetadata, List[str], List[str]]:
        """Process PowerPoint files (.pptx, .ppt)"""
        warnings = []
        text_pages = []
        images_base64 = []

        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            # Extract text from each slide (one slide = one page)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract images
                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            img_base64 = base64.b64encode(image_bytes).decode('utf-8')
                            images_base64.append(img_base64)
                        except Exception:
                            pass
                
                page_content = f"[Slide {slide_num}]\n" + '\n'.join(slide_text)
                text_pages.append(page_content if slide_text else f"[Slide {slide_num} - No text]")
            
            if not text_pages:
                text_pages = ["[No slides found]"]
            
            metadata = DocumentMetadata(
                filename=Path(file_path).name,
                file_size=os.path.getsize(file_path),
                page_count=len(text_pages),
                image_count=len(images_base64),
                is_legible=True,
                legibility_score=999.0,
                format="PPTX" if file_path.endswith('.pptx') else "PPT",
                warnings=warnings
            )
            
            return metadata, text_pages, images_base64
            
        except ImportError:
            raise ValueError("python-pptx not installed. Run: pip install python-pptx")
        except Exception as e:
            raise ValueError(f"Error processing PowerPoint: {str(e)}")

    def _process_excel(self, file_path: str) -> Tuple[DocumentMetadata, List[str], List[str]]:
        """Process Excel files (.xlsx, .xls)"""
        warnings = []
        text_pages = []
        images_base64 = []

        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, data_only=True)
            
            # Process each sheet as a "page"
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_text = [f"[Sheet: {sheet_name}]"]
                
                # Extract cell values
                for row in ws.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_text):  # Skip empty rows
                        sheet_text.append(" | ".join(row_text))
                
                text_pages.append('\n'.join(sheet_text))
            
            if not text_pages:
                text_pages = ["[No sheets found]"]
            
            metadata = DocumentMetadata(
                filename=Path(file_path).name,
                file_size=os.path.getsize(file_path),
                page_count=len(text_pages),
                image_count=0,  # Excel images are complex to extract
                is_legible=True,
                legibility_score=999.0,
                format="XLSX" if file_path.endswith('.xlsx') else "XLS",
                warnings=warnings
            )
            
            return metadata, text_pages, images_base64
            
        except ImportError:
            raise ValueError("openpyxl not installed. Run: pip install openpyxl")
        except Exception as e:
            raise ValueError(f"Error processing Excel: {str(e)}")

    def _process_text(self, file_path: str) -> Tuple[DocumentMetadata, List[str], List[str]]:
        """Process plain text files (.txt, .md, .csv)"""
        warnings = []
        text_pages = []
        images_base64 = []

        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # Split into pages (every 2000 characters or by line breaks)
            if len(content) <= 2000:
                text_pages = [content]
            else:
                # Split by approximate page size
                page_size = 2000
                text_pages = [content[i:i+page_size] for i in range(0, len(content), page_size)]
            
            if not text_pages or not content.strip():
                text_pages = ["[Empty file]"]
            
            file_ext = Path(file_path).suffix.upper().replace('.', '')
            
            metadata = DocumentMetadata(
                filename=Path(file_path).name,
                file_size=os.path.getsize(file_path),
                page_count=len(text_pages),
                image_count=0,
                is_legible=True,
                legibility_score=999.0,
                format=file_ext,
                warnings=warnings
            )
            
            return metadata, text_pages, images_base64
            
        except Exception as e:
            raise ValueError(f"Error processing text file: {str(e)}")

    def _check_image_legibility(self, img: Image.Image) -> float:
        """
        Check image legibility using Laplacian variance (blur detection)
        Higher values indicate sharper images
        """
        try:
            # Convert to grayscale
            gray = img.convert('L')
            gray_array = np.array(gray)

            # Resize if too large (for performance)
            max_dimension = 1000
            if max(gray_array.shape) > max_dimension:
                scale = max_dimension / max(gray_array.shape)
                new_size = (int(gray_array.shape[1] * scale), int(gray_array.shape[0] * scale))
                gray = gray.resize(new_size, Image.Resampling.LANCZOS)
                gray_array = np.array(gray)

            # Calculate Laplacian variance
            laplacian = cv2.Laplacian(gray_array, cv2.CV_64F)
            variance = laplacian.var()

            return float(variance)

        except Exception as e:
            # If check fails, assume legible
            return self.MIN_LEGIBILITY_SCORE

    def extract_document_info(self, text_pages: List[str]) -> str:
        """
        Combine text from all pages into a single string for analysis
        """
        combined_text = "\n\n--- Page Break ---\n\n".join(
            [f"Page {i+1}:\n{text}" for i, text in enumerate(text_pages)]
        )
        return combined_text

    def prepare_claude_content(
        self,
        text_pages: List[str],
        images_base64: List[str],
        max_images: int = 20  # Claude has limits on image count
    ) -> List[Dict[str, Any]]:
        """
        Prepare content blocks for Claude API (text + images)
        """
        content_blocks = []

        # Add text content
        if text_pages:
            combined_text = self.extract_document_info(text_pages)
            content_blocks.append({
                "type": "text",
                "text": f"Document Text Content:\n\n{combined_text}"
            })

        # Add images (limit to avoid API issues)
        for idx, img_base64 in enumerate(images_base64[:max_images]):
            # Detect media type from base64 data
            try:
                img_bytes = base64.b64decode(img_base64)
                media_type = self._detect_image_format(img_bytes)
            except Exception:
                media_type = "image/png"  # Fallback if decoding fails

            content_blocks.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": media_type,
                    "data": img_base64
                }
            })

        if len(images_base64) > max_images:
            content_blocks.append({
                "type": "text",
                "text": f"\n\nNote: This document contains {len(images_base64)} images, but only the first {max_images} are shown due to API limitations."
            })

        return content_blocks

    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get basic file information without full processing"""
        return {
            "filename": Path(file_path).name,
            "file_size": os.path.getsize(file_path),
            "format": Path(file_path).suffix.upper().replace('.', ''),
            "path": file_path
        }

    def extract_page_images(self, file_path: str) -> List[str]:
        """
        Extract page renderings as base64 images for preview in UI
        Returns list of base64 encoded PNG images (one per page)
        """
        file_ext = Path(file_path).suffix.lower()
        page_images = []

        try:
            if file_ext == '.pdf':
                doc = fitz.open(file_path)
                for page in doc:
                    # Render at high quality
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img_bytes = pix.tobytes("png")
                    img_base64 = base64.b64encode(img_bytes).decode('utf-8')
                    page_images.append(img_base64)
                doc.close()
                
            elif file_ext == '.docx':
                # Render Word document pages as images
                page_images = self._render_word_pages(file_path)
                
            elif file_ext == '.pptx':
                # Render PowerPoint slides as images
                page_images = self._render_ppt_slides(file_path)
                
            elif file_ext == '.xlsx':
                # Render Excel sheets as images
                page_images = self._render_excel_sheets(file_path)
                
            elif file_ext in {'.txt', '.md', '.csv'}:
                # Render text files as images
                page_images = self._render_text_pages(file_path)
                
            else:
                # For image files, just encode the image
                with open(file_path, 'rb') as f:
                    image_bytes = f.read()
                img_base64 = base64.b64encode(image_bytes).decode('utf-8')
                page_images.append(img_base64)

        except Exception as e:
            print(f"Error extracting page images: {str(e)}")
            # Return a fallback error image
            page_images = [self._create_error_image(str(e))]

        return page_images

    def _render_word_pages(self, file_path: str) -> List[str]:
        """Render Word document pages as images"""
        from PIL import Image, ImageDraw, ImageFont
        from docx import Document
        
        print(f"[DEBUG] DEBUG: Starting Word rendering for {file_path}")
        page_images = []
        
        try:
            doc = Document(file_path)
            print(f"[DEBUG] DEBUG: Document loaded, {len(doc.paragraphs)} paragraphs found")
            
            # Group paragraphs into pages (roughly)
            current_page_text = []
            line_count = 0
            max_lines_per_page = 40
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # Estimate lines (rough)
                    estimated_lines = max(1, len(text) // 80)
                    
                    if line_count + estimated_lines > max_lines_per_page and current_page_text:
                        # Render current page
                        img_base64 = self._create_text_image(current_page_text, f"Word Document - Page {len(page_images) + 1}")
                        page_images.append(img_base64)
                        current_page_text = []
                        line_count = 0
                    
                    current_page_text.append(text)
                    line_count += estimated_lines
            
            # Render last page
            if current_page_text:
                img_base64 = self._create_text_image(current_page_text, f"Word Document - Page {len(page_images) + 1}")
                page_images.append(img_base64)
            
            print(f"[DEBUG] DEBUG: Created {len(page_images)} page images")
            
            # If no content, create placeholder
            if not page_images:
                print(f"[DEBUG] DEBUG: No pages created, adding placeholder")
                page_images.append(self._create_text_image(["[Empty Document]"], "Word Document"))
                
        except Exception as e:
            print(f"[ERROR] ERROR: Error rendering Word pages: {str(e)}")
            import traceback
            traceback.print_exc()
            page_images = [self._create_error_image(f"Word rendering error: {str(e)}")]
        
        return page_images

    def _render_ppt_slides(self, file_path: str) -> List[str]:
        """Render PowerPoint slides as images"""
        from pptx import Presentation
        
        page_images = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract all text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Create image for this slide
                img_base64 = self._create_text_image(
                    slide_text if slide_text else ["[No text on slide]"],
                    f"Slide {slide_num} of {len(prs.slides)}"
                )
                page_images.append(img_base64)
            
            if not page_images:
                page_images.append(self._create_text_image(["[Empty Presentation]"], "PowerPoint"))
                
        except Exception as e:
            print(f"Error rendering PPT slides: {str(e)}")
            page_images = [self._create_error_image(f"PowerPoint rendering error: {str(e)}")]
        
        return page_images

    def _render_excel_sheets(self, file_path: str) -> List[str]:
        """Render Excel sheets as images"""
        from openpyxl import load_workbook
        
        page_images = []
        
        try:
            wb = load_workbook(file_path, data_only=True)
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_lines = [f"Sheet: {sheet_name}", "=" * 50]
                
                # Get data (limit to first 30 rows for readability)
                for idx, row in enumerate(ws.iter_rows(values_only=True, max_row=30), 1):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_text):
                        sheet_lines.append(f"{idx:3d} | " + " | ".join(row_text[:10]))  # First 10 columns
                
                if len(list(ws.iter_rows())) > 30:
                    sheet_lines.append(f"... ({len(list(ws.iter_rows()))} total rows)")
                
                img_base64 = self._create_text_image(sheet_lines, f"Excel - {sheet_name}")
                page_images.append(img_base64)
            
            if not page_images:
                page_images.append(self._create_text_image(["[Empty Workbook]"], "Excel"))
                
        except Exception as e:
            print(f"Error rendering Excel sheets: {str(e)}")
            page_images = [self._create_error_image(f"Excel rendering error: {str(e)}")]
        
        return page_images

    def _render_text_pages(self, file_path: str) -> List[str]:
        """Render text files as images"""
        page_images = []
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # Split into pages (40 lines per page)
            lines = content.split('\n')
            lines_per_page = 40
            
            for i in range(0, len(lines), lines_per_page):
                page_lines = lines[i:i+lines_per_page]
                img_base64 = self._create_text_image(
                    page_lines,
                    f"Text File - Page {i//lines_per_page + 1}"
                )
                page_images.append(img_base64)
            
            if not page_images:
                page_images.append(self._create_text_image(["[Empty File]"], "Text File"))
                
        except Exception as e:
            print(f"Error rendering text pages: {str(e)}")
            page_images = [self._create_error_image(f"Text rendering error: {str(e)}")]
        
        return page_images

    def _create_text_image(self, lines: List[str], header: str = "Document") -> str:
        """Create an image from text lines with better formatting"""
        from PIL import Image, ImageDraw, ImageFont
        
        print(f"[DEBUG] DEBUG: Creating text image with {len(lines)} lines, header: {header}")
        
        # Image dimensions
        width = 1200
        height = 1600
        padding = 60
        line_height = 32
        
        # Create image with white background
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)
        
        try:
            # Try to use a nice font
            title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 28)
            text_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
        except:
            # Fallback to default font
            print(f"[WARNING] DEBUG: Using default font")
            title_font = ImageFont.load_default()
            text_font = ImageFont.load_default()
        
        # Draw header background
        draw.rectangle([0, 0, width, 100], fill='#002E6D')  # Hitachi navy
        draw.text((padding, 35), header, fill='white', font=title_font)
        
        # Draw border
        draw.rectangle([0, 0, width-1, height-1], outline='#002E6D', width=3)
        
        # Draw text content
        y_position = 120
        
        for line in lines:
            if y_position > height - padding:
                # Add "..." if content overflows
                draw.text((padding, y_position), "...", fill='#666666', font=text_font)
                break
            
            # Wrap long lines
            if len(line) > 90:
                line = line[:87] + "..."
            
            draw.text((padding, y_position), line, fill='#333333', font=text_font)
            y_position += line_height
        
        # Convert to base64
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        
        return img_base64

    def _create_error_image(self, error_message: str) -> str:
        """Create an error image"""
        from PIL import Image, ImageDraw, ImageFont
        
        width = 1200
        height = 800
        
        img = Image.new('RGB', (width, height), color='#FFF5F5')
        draw = ImageDraw.Draw(img)
        
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
        except:
            font = ImageFont.load_default()
        
        # Draw error message
        draw.rectangle([0, 0, width-1, height-1], outline='#DC2626', width=5)
        draw.text((100, 300), "[WARNING] Error rendering document", fill='#DC2626', font=font)
        draw.text((100, 400), error_message[:100], fill='#666666', font=font)
        
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        
        return img_base64

